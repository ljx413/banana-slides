import streamlit as st
from pptx import Presentation
import io

st.set_page_config(page_title="Banana Slides", layout="centered")
st.title("🍌 Banana Slides - 一键 PPT 生成器")

uploaded = st.file_uploader("上传 PDF / DOCX / TXT", type=["pdf", "docx", "txt"])
if uploaded and st.button("生成 PPT"):
    with st.spinner("正在生成，请稍候..."):
        # ===== 这里直接生成 =====
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Hello Banana Slides"
        slide.placeholders[1].text = f"你上传了：{uploaded.name}"
        # 保存到内存
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
    st.download_button("📥 下载 PPT", buffer, file_name="banana_slides.pptx")

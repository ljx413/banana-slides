import os, uuid, shutil
from fastapi import FastAPI, UploadFile, File
from fastapi.responses import FileResponse
from pathlib import Path

app = FastAPI(title="BananaSlides")

# 临时目录
TMP = Path("/tmp")
TMP.mkdir(exist_ok=True)

@app.post("/generate")
def generate_ppt(upload: UploadFile = File(...)):
    uid   = uuid.uuid4().hex
    in_file  = TMP / f"{uid}_{upload.filename}"
    out_file = TMP / f"{uid}.pptx"

    # 保存上传文件
    with in_file.open("wb") as f:
        shutil.copyfileobj(upload.file, f)

    # ===== 这里调你已有的「services」伪代码，先直接复制一份做演示 =====
    # 如果你后面写好了 services.generate_ppt(str(in_file), str(out_file))
    # 把下面两行换成真实调用即可
    import python_pptx
    from pptx import Presentation
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0]).shapes.title.text = "Hello Banana Slides"
    prs.save(out_file)
    # =========================================================

    return FileResponse(out_file, filename="banana_slides.pptx")
